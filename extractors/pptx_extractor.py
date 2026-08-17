#!/usr/bin/env python3
"""
.pptx extractor (Office Open XML PowerPoint presentations)

Responsibilities
---------------
1) Full-text extraction for .pptx files: per slide, in slide order —
   text-frame/shape text (titles, bullets, text boxes, including text
   inside grouped shapes), table cell text, and speaker notes.
2) Lightweight metadata extraction when `metadata_only=True`.

Design notes
------------
- `python-pptx` parses the OOXML package directly.
- Embedded images are not OCR'd,
  embedded objects/chart data are not parsed, and slide-master/template
  boilerplate is not deduplicated — only the shapes actually placed on
  each slide (and its notes slide) are read.
- A single unreadable slide, shape, table, or notes slide is recorded in
  `pptx_parts_failed`/`pptx_part_failure_reasons` and skipped rather than
  aborting the whole presentation — mirrors eml_extractor.py's per-part
  failure contract. `extract_pptx()` only raises ExtractionError when the
  file itself can't be read or isn't parseable as a .pptx package at all.
- The binary-content guard (text_extractor.looks_like_binary) does not
  apply here: discovery.py only runs it for TEXT_EXTENSIONS, and the text
  this module returns is already decoded, never raw file bytes.
"""

from __future__ import annotations

import os
from typing import Any, Dict, List, Tuple

from .errors import ExtractionError

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # pragma: no cover - dependency is pinned/required
    Presentation = None
    MSO_SHAPE_TYPE = None


def _iter_leaf_shapes(shapes):
    """Recurse into group shapes so text nested inside a group is not lost."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def _extract_shape_text(shape) -> List[str]:
    lines: List[str] = []
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        text = shape.text_frame.text
        if text and text.strip():
            lines.append(text.strip())
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
    return lines


def _extract_slide_text(slide, slide_num: int) -> Tuple[List[str], int, List[str]]:
    """Returns (lines, parts_failed, failure_reasons) for one slide."""
    lines: List[str] = []
    parts_failed = 0
    failure_reasons: List[str] = []

    try:
        leaf_shapes = list(_iter_leaf_shapes(slide.shapes))
    except Exception as exc:
        return [], 1, [f"slide {slide_num}: {type(exc).__name__}: {exc}"]

    for shape in leaf_shapes:
        try:
            lines.extend(_extract_shape_text(shape))
        except Exception as exc:
            parts_failed += 1
            failure_reasons.append(
                f"slide {slide_num} shape: {type(exc).__name__}: {exc}"
            )

    try:
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text and notes_text.strip():
                lines.append(f"Notes: {notes_text.strip()}")
    except Exception as exc:
        parts_failed += 1
        failure_reasons.append(f"slide {slide_num} notes: {type(exc).__name__}: {exc}")

    return lines, parts_failed, failure_reasons


def _normalize_pptx_date(raw) -> "str | None":
    if raw is None:
        return None
    try:
        return raw.isoformat()
    except Exception:
        return str(raw)


def _extract_pptx_metadata(prs) -> Dict[str, Any]:
    meta: Dict[str, Any] = {}
    try:
        props = prs.core_properties
        mapping = {
            "title": props.title,
            "author": props.author,
            "subject": props.subject,
            "created": _normalize_pptx_date(props.created),
            "modified": _normalize_pptx_date(props.modified),
        }
        for k, v in mapping.items():
            if v not in (None, "", "None"):
                meta[k] = v
    except Exception:
        pass
    return meta


def _open_pptx(path: str):
    try:
        return Presentation(path)
    except Exception as exc:
        raise ExtractionError.from_exception(exc) from exc


def extract_pptx(path: str, metadata_only: bool = False, return_details: bool = False):
    """
    Public API used by discovery.py.

    - metadata_only=True  -> Dict[str, Any] of core document properties
      (title/author/subject/created/modified, when present) plus
      slide_count.
    - metadata_only=False -> str: for each slide in order, one "Slide N"
      label line, then text-frame/shape text (titles, bullets, text
      boxes — including text nested in grouped shapes), then one line
      per table row (" | "-joined cells), then a "Notes: ..." line if
      the slide has speaker notes. Slides that contribute no text are
      omitted entirely.
    - return_details=True (text mode only) additionally returns a dict
      with pptx_parts_failed / pptx_part_failure_reasons, present only
      when nonzero — mirrors extract_eml()'s return_details contract.

    Raises ExtractionError only when the file itself can't be read or
    isn't parseable as a .pptx package at all. Individual unreadable
    slides/shapes/tables/notes are degraded, not fatal — see module
    docstring.
    """
    if not isinstance(path, str):
        raise ExtractionError(f"TypeError: path must be str, got {type(path).__name__}")
    if not path or not path.strip():
        raise ExtractionError("ValueError: path cannot be empty")
    if not os.path.isfile(path):
        raise ExtractionError("FileNotFoundError: file does not exist")
    if Presentation is None:
        raise ExtractionError("RuntimeError: python-pptx is not installed")

    prs = _open_pptx(path)

    if metadata_only:
        meta = _extract_pptx_metadata(prs)
        try:
            meta["slide_count"] = len(prs.slides)
        except Exception:
            pass
        return meta

    blocks: List[str] = []
    total_parts_failed = 0
    all_failure_reasons: List[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            lines, parts_failed, failure_reasons = _extract_slide_text(slide, slide_num)
        except Exception as exc:
            lines, parts_failed, failure_reasons = (
                [],
                1,
                [f"slide {slide_num}: {type(exc).__name__}: {exc}"],
            )
        total_parts_failed += parts_failed
        all_failure_reasons.extend(failure_reasons)
        if lines:
            blocks.append(f"Slide {slide_num}")
            blocks.extend(lines)

    text = "\n".join(blocks)

    if not return_details:
        return text

    details: Dict[str, Any] = {}
    if total_parts_failed:
        details["pptx_parts_failed"] = total_parts_failed
        details["pptx_part_failure_reasons"] = all_failure_reasons
    return text, details
