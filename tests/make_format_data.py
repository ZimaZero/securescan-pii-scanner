#!/usr/bin/env python3
"""
Generates a format-coverage corpus at tests/format_data/ plus a verification
manifest at tests/format_data/manifest.json.

Every supported extraction path (docx, xlsx, text-layer pdf, image-only pdf,
png, csv, json) gets at least one file. A fixed, pre-validated PII set is
planted at documented locations — including locations expected to expose
extractor gaps — so the
test suite that reads this manifest (tests/test_format_coverage.py) can tell
a real detection gap apart from a real extraction gap.

Regenerate with:
    docker compose run --rm securescan-cpu python tests/make_format_data.py
"""

import json
import os
import sys

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

from detectors.detectors import luhn_check  # noqa: E402
from detectors.keyword_detector import validate_sin  # noqa: E402
from detectors.health_card_detector import ohip_valid, bc_phn_valid  # noqa: E402

OUT_DIR = os.path.join(PROJECT_ROOT, "tests", "format_data")
MANIFEST_PATH = os.path.join(OUT_DIR, "manifest.json")

FONT_REGULAR = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"

# ---------------------------------------------------------------------------
# Fixed PII set — verify every value with the project validators below before
# writing a fixture. Hand-picked identifiers must pass the real checksum.
# ---------------------------------------------------------------------------
SIN = "132-677-360"
ON_HEALTH_CARD = "9876543217"
BC_HEALTH_CARD = "9123456780"
CA_PASSPORT = "AB123456"          # format-only, no checksum; needs "Passport:" keyword
ON_DRIVERS_LICENSE = "A1234-56789-01234"  # format-only; needs a DL keyword
CREDIT_CARD = "4111 1111 1111 1111"
EMAIL = "format.test@example.com"
PHONE = "416-555-0199"
POSTAL_CODE = "T2X 1V4"           # needs a postal keyword nearby


def verify_planted_values():
    checks = [
        ("SIN", validate_sin(SIN)),
        ("ON health card (OHIP)", ohip_valid(ON_HEALTH_CARD)),
        ("BC health card (PHN)", bc_phn_valid(BC_HEALTH_CARD)),
        ("Credit card (Luhn)", luhn_check(CREDIT_CARD)),
    ]
    failed = [name for name, ok in checks if not ok]
    if failed:
        raise SystemExit(
            "Refusing to generate corpus — these planted values FAILED their "
            f"own validator: {failed}. Fix the value, don't fix the test."
        )
    for name, _ in checks:
        print(f"[✓] {name} verified valid")
    # CA_PASSPORT / ON_DRIVERS_LICENSE / EMAIL / PHONE / POSTAL_CODE have no
    # checksum — format + keyword is the whole contract, verified later by
    # actually running detection over the generated files.


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def build_docx_full(path):
    from docx import Document

    doc = Document()
    doc.add_heading("Format Coverage Test Document", level=1)
    doc.add_paragraph(f"Employee SIN on file: {SIN}")
    doc.add_paragraph("This paragraph intentionally contains no other PII.")

    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Field"
    table.rows[0].cells[1].text = "Value"
    table.rows[1].cells[0].text = "Card on file"
    table.rows[1].cells[1].text = CREDIT_CARD

    section = doc.sections[0]
    section.header.paragraphs[0].text = f"CONFIDENTIAL — contact {EMAIL}"
    section.footer.paragraphs[0].text = f"Questions? Call {PHONE}"

    doc.save(path)


def build_docx_no_pii(path):
    from docx import Document

    doc = Document()
    doc.add_heading("Quarterly Planning Notes", level=1)
    doc.add_paragraph(
        "The team reviewed the roadmap and agreed on three priorities for "
        "next quarter: platform reliability, documentation, and onboarding."
    )
    doc.add_paragraph(
        "No action items require follow-up outside the usual weekly sync."
    )
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Topic"
    table.rows[0].cells[1].text = "Owner"
    table.rows[1].cells[0].text = "Documentation"
    table.rows[1].cells[1].text = "Team Lead"
    doc.save(path)


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

def build_xlsx_multisheet(path):
    from openpyxl import Workbook
    from openpyxl.comments import Comment

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Records"
    ws1["A1"] = "Field"
    ws1["B1"] = "Value"
    ws1["A2"] = "SIN"
    ws1["B2"] = SIN
    ws1["A3"] = "Credit Card"
    ws1["B3"] = CREDIT_CARD
    ws1["B3"].number_format = "@"  # explicitly text-formatted, not numeric

    # Cell comment: a value that lives ONLY in a comment, never in a cell's
    # own text. openpyxl's read_only load (used by xlsx_extractor.py) does
    # not surface comments — this is expected to be a real extraction gap,
    # not something the manifest should quietly work around.
    ws1["A4"] = "Internal note"
    ws1["A4"].comment = Comment(f"Health card on file: {ON_HEALTH_CARD}", "generator")

    ws2 = wb.create_sheet("Contact")
    ws2["A1"] = "Email"
    ws2["B1"] = EMAIL
    ws2["A2"] = "Phone"
    ws2["B2"] = PHONE
    ws2["A3"] = "Postal Code"
    ws2["B3"] = POSTAL_CODE

    wb.save(path)


# ---------------------------------------------------------------------------
# PDF — text layer (reportlab)
# ---------------------------------------------------------------------------

def build_pdf_text_layer(path):
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter

    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(path, pagesize=letter)
    story = [
        Paragraph("Format Coverage — Text-Layer PDF", styles["Title"]),
        Spacer(1, 12),
        Paragraph(f"SIN on file: {SIN}", styles["Normal"]),
        Paragraph(f"Contact email: {EMAIL}", styles["Normal"]),
        Spacer(1, 12),
    ]
    table_data = [["Field", "Value"], ["Card on file", CREDIT_CARD], ["Phone", PHONE]]
    table = Table(table_data)
    table.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
        ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
    ]))
    story.append(table)
    doc.build(story)


def build_pdf_no_pii(path):
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.pagesizes import letter

    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(path, pagesize=letter)
    story = [
        Paragraph("Quarterly Planning Notes", styles["Title"]),
        Spacer(1, 12),
        Paragraph(
            "This document summarizes the team's planning discussion. No "
            "customer or personal data is referenced anywhere in this file.",
            styles["Normal"],
        ),
        Paragraph(
            "Next steps will be tracked in the usual project board.",
            styles["Normal"],
        ),
    ]
    doc.build(story)


# ---------------------------------------------------------------------------
# PDF — image-only ("scanned"), 2 pages, PII split across pages
# ---------------------------------------------------------------------------

def _render_page_image(lines, font_size=44):
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (1700, 2200), color="white")
    draw = ImageDraw.Draw(img)
    font = ImageFont.truetype(FONT_REGULAR, font_size)
    y = font_size
    for line in lines:
        draw.text((80, y), line, fill="black", font=font)
        y += int(font_size * 1.8)
    return img


def build_pdf_scanned_2page(path):
    import fitz

    page1_lines = [
        "SCANNED RECORD — PAGE 1 OF 2",
        "",
        f"Passport: {CA_PASSPORT}",
        f"Driver's Licence: {ON_DRIVERS_LICENSE}",
    ]
    page2_lines = [
        "SCANNED RECORD — PAGE 2 OF 2",
        "",
        f"Ontario Health Card: {ON_HEALTH_CARD}",
        f"BC Health Card (PHN): {BC_HEALTH_CARD}",
    ]

    tmp1 = os.path.join(OUT_DIR, "_scan_p1.jpg")
    tmp2 = os.path.join(OUT_DIR, "_scan_p2.jpg")
    _render_page_image(page1_lines).convert("RGB").save(tmp1, "JPEG", quality=85)
    _render_page_image(page2_lines).convert("RGB").save(tmp2, "JPEG", quality=85)

    doc = fitz.open()
    for tmp in (tmp1, tmp2):
        page = doc.new_page(width=612, height=792)
        page.insert_image(fitz.Rect(0, 0, 612, 792), filename=tmp)
    doc.save(path)
    doc.close()
    os.remove(tmp1)
    os.remove(tmp2)


# ---------------------------------------------------------------------------
# PNG images — clean / small-font / rotated (OCR robustness)
# ---------------------------------------------------------------------------

def build_png(path, font_size, rotate_deg=0):
    from PIL import Image, ImageDraw, ImageFont

    lines = [
        f"SIN: {SIN}",
        f"Email: {EMAIL}",
        f"Phone: {PHONE}",
    ]
    canvas_w, canvas_h = 1000, 260
    img = Image.new("RGB", (canvas_w, canvas_h), color="white")
    draw = ImageDraw.Draw(img)
    font = ImageFont.truetype(FONT_REGULAR, font_size)
    y = 20
    for line in lines:
        draw.text((40, y), line, fill="black", font=font)
        y += int(font_size * 1.7)

    if rotate_deg:
        img = img.rotate(rotate_deg, expand=True, fillcolor="white")

    img.save(path)


def build_png_exif_rotated(path, font_size=32):
    """Upright content, but the pixel data is pre-rotated 270deg clockwise
    (== PIL's img.rotate(90), a counter-clockwise convention) and tagged
    with EXIF Orientation=6 — the pairing a real EXIF-honoring viewer (and
    now image_extractor.py's step-0 exif_transpose) restores to upright.
    Probes the OCR ladder's cheapest rung: this should score identically to
    image_clean.png WITHOUT ever reaching OSD/brute-force/deskew, since
    exif_transpose() alone fixes it before any of those run."""
    from PIL import Image, ImageDraw, ImageFont

    lines = [
        f"SIN: {SIN}",
        f"Email: {EMAIL}",
        f"Phone: {PHONE}",
    ]
    canvas_w, canvas_h = 1000, 260
    img = Image.new("RGB", (canvas_w, canvas_h), color="white")
    draw = ImageDraw.Draw(img)
    font = ImageFont.truetype(FONT_REGULAR, font_size)
    y = 20
    for line in lines:
        draw.text((40, y), line, fill="black", font=font)
        y += int(font_size * 1.7)

    rotated = img.rotate(90, expand=True, fillcolor="white")
    exif = rotated.getexif()
    exif[0x0112] = 6  # Orientation tag
    rotated.save(path, exif=exif)


# ---------------------------------------------------------------------------
# CSV / JSON
# ---------------------------------------------------------------------------

def build_csv(path):
    import csv as csv_mod

    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv_mod.writer(f)
        writer.writerow(["id", "name", "sin", "email", "phone", "postal_code"])
        # Row labels deliberately avoid "test"/"sample"/etc. — those are real
        # placeholder-suppression trigger words (see keyword_detector.py's
        # TEST_INDICATORS) and would legitimately suppress the PII in this
        # same row, which isn't what this fixture is testing.
        writer.writerow([1, "Record One", SIN, EMAIL, PHONE, f"Postal Code {POSTAL_CODE}"])
        writer.writerow([2, "Record Two", "", "", "", ""])


# ---------------------------------------------------------------------------
# EML — multipart email: plain-text body, HTML body, one attachment name
# ---------------------------------------------------------------------------

def build_eml(path):
    from email.message import EmailMessage

    msg = EmailMessage()
    msg["From"] = "records@example.org"
    msg["To"] = "case.worker@example.org"
    msg["Subject"] = "Account summary"
    msg["Date"] = "Thu, 04 Jun 2026 10:00:00 -0000"
    msg.set_content(f"SIN on file: {SIN}\nContact phone: {PHONE}\n")
    msg.add_alternative(
        "<html><body>"
        f"<p>Card on file: {CREDIT_CARD}</p>"
        f"<p>Postal Code: {POSTAL_CODE}</p>"
        "</body></html>",
        subtype="html",
    )
    msg.add_attachment(
        b"%PDF-1.4 fake invoice body",
        maintype="application",
        subtype="pdf",
        filename="invoice.pdf",
    )
    with open(path, "wb") as f:
        f.write(msg.as_bytes())


# ---------------------------------------------------------------------------
# PPTX — title slide, table slide, speaker notes
# ---------------------------------------------------------------------------

def build_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text_frame.text = "Case Summary"
    body = slide1.placeholders[1]
    body.text_frame.text = f"SIN on file: {SIN}"
    p2 = body.text_frame.add_paragraph()
    p2.text = f"Contact phone: {PHONE}"
    slide1.notes_slide.notes_text_frame.text = f"Confidential — Postal Code: {POSTAL_CODE}"

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = width = height = Inches(1)
    table_shape = slide2.shapes.add_table(2, 2, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "Field"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Card on file"
    table.cell(1, 1).text = CREDIT_CARD

    prs.save(path)


def build_json(path):
    data = [
        {
            "id": 1,
            "sin": SIN,
            "email": EMAIL,
            "phone": PHONE,
            "note": f"Postal Code: {POSTAL_CODE}",
        },
        {"id": 2, "note": "no PII in this record"},
    ]
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)


# ---------------------------------------------------------------------------
# Manifest
# ---------------------------------------------------------------------------

def finding(taxonomy, value, band, location, note=None, known_fail=False):
    d = {"taxonomy": taxonomy, "value": value, "band": band, "location": location}
    if note:
        d["note"] = note
    if known_fail:
        # Documented, accepted limitation: still checked and reported by the
        # test suite (as XFAIL/XPASS), but does NOT count toward the file's
        # overall pass/fail — unlike an ordinary miss, this one is expected.
        d["known_fail"] = True
    return d


def build_manifest():
    files = [
        {
            "path": "tests/format_data/docx_full.docx",
            "format": "docx",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "body paragraph"),
                finding("identifier.financial.credit_card", CREDIT_CARD, "HIGH", "table cell"),
                finding(
                    "contact.email", EMAIL, "MEDIUM", "header",
                    note="docx_extractor.py's _extract_docx_text() only reads "
                         "doc.paragraphs + doc.tables, not section headers/footers "
                         "— expected extraction gap, not a detection gap.",
                ),
                finding(
                    "contact.phone", PHONE, "MEDIUM", "footer",
                    note="Same header/footer extraction gap as the email above.",
                ),
            ],
        },
        {
            "path": "tests/format_data/docx_no_pii.docx",
            "format": "docx",
            "negative_control": True,
            "expected": [],
        },
        {
            "path": "tests/format_data/xlsx_multisheet.xlsx",
            "format": "xlsx",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "Sheet 'Records', cell B2"),
                finding("identifier.financial.credit_card", CREDIT_CARD, "HIGH",
                        "Sheet 'Records', cell B3 (text-formatted)"),
                finding(
                    "identifier.government.health_card_on", ON_HEALTH_CARD, "HIGH",
                    "Sheet 'Records', cell A4 comment",
                    note="xlsx_extractor.py loads with read_only=True, which does not "
                         "expose openpyxl cell comments via iter_rows() — expected "
                         "extraction gap, not a detection gap.",
                ),
                finding("contact.email", EMAIL, "MEDIUM", "Sheet 'Contact', cell B1"),
                finding("contact.phone", PHONE, "MEDIUM", "Sheet 'Contact', cell B2"),
                finding("contact.address.postal_code", POSTAL_CODE, "MEDIUM",
                        "Sheet 'Contact', cell B3 (label 'Postal Code' in A3)"),
            ],
        },
        {
            "path": "tests/format_data/pdf_text_layer.pdf",
            "format": "pdf",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "paragraph"),
                finding("contact.email", EMAIL, "MEDIUM", "paragraph"),
                finding("identifier.financial.credit_card", CREDIT_CARD, "HIGH", "table cell"),
                finding("contact.phone", PHONE, "MEDIUM", "table cell"),
            ],
        },
        {
            "path": "tests/format_data/pdf_no_pii.pdf",
            "format": "pdf",
            "negative_control": True,
            "expected": [],
        },
        {
            "path": "tests/format_data/pdf_scanned_2page.pdf",
            "format": "pdf",
            "negative_control": False,
            "expected": [
                finding("identifier.government.passport_ca", CA_PASSPORT, "HIGH",
                        "page 1, OCR"),
                finding("identifier.government.drivers_license_on", ON_DRIVERS_LICENSE,
                        "HIGH", "page 1, OCR"),
                finding("identifier.government.health_card_on", ON_HEALTH_CARD, "HIGH",
                        "page 2, OCR"),
                finding("identifier.government.health_card_bc", BC_HEALTH_CARD, "HIGH",
                        "page 2, OCR"),
            ],
        },
        {
            "path": "tests/format_data/image_clean.png",
            "format": "png",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "large clear text, OCR"),
                finding("contact.email", EMAIL, "MEDIUM", "large clear text, OCR"),
                finding("contact.phone", PHONE, "MEDIUM", "large clear text, OCR"),
            ],
        },
        {
            "path": "tests/format_data/image_small.png",
            "format": "png",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "~14pt text, OCR",
                        note="OCR robustness probe — small font size."),
                finding("contact.email", EMAIL, "MEDIUM", "~14pt text, OCR",
                        note="OCR robustness probe — small font size."),
                finding("contact.phone", PHONE, "MEDIUM", "~14pt text, OCR",
                        note="OCR robustness probe — small font size."),
            ],
        },
        {
            "path": "tests/format_data/image_rotated.png",
            "format": "png",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "rotated 5deg, OCR",
                        note="OCR robustness probe — rotation."),
                finding("contact.email", EMAIL, "MEDIUM", "rotated 5deg, OCR",
                        note="OCR robustness probe — rotation."),
                finding("contact.phone", PHONE, "MEDIUM", "rotated 5deg, OCR",
                        note="Was a documented known_fail: at 5deg rotation "
                             "Tesseract misread the leading '4' in "
                             "'416-555-0199' as 'A', breaking the phone regex, "
                             "even though Tesseract's own confidence stayed "
                             "high (a wrong-but-confident read). Fixed by "
                             "image_extractor.py's orientation-resilience "
                             "ladder: OSD + brute-force cardinal correction for "
                             "90/180/270 rotations, plus a fine-angle deskew "
                             "rung (always attempted, judged on its own "
                             "geometric signal rather than OCR confidence, "
                             "since confidence can't detect this failure mode) "
                             "that now catches this sub-cardinal skew."),
            ],
        },
        {
            "path": "tests/format_data/image_rotated_180.png",
            "format": "png",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "rotated 180deg, OCR",
                        note="OCR robustness probe — full upside-down orientation, "
                             "the realistic case for a photographed ID card. "
                             "Corrected by image_extractor.py's OSD step."),
                finding("contact.email", EMAIL, "MEDIUM", "rotated 180deg, OCR",
                        note="OCR robustness probe — rotation."),
                finding("contact.phone", PHONE, "MEDIUM", "rotated 180deg, OCR",
                        note="OCR robustness probe — rotation."),
            ],
        },
        {
            "path": "tests/format_data/image_rotated_270.png",
            "format": "png",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "rotated 270deg, OCR",
                        note="OCR robustness probe — sideways orientation, the "
                             "realistic case for a photographed ID card. "
                             "Corrected by image_extractor.py's OSD step."),
                finding("contact.email", EMAIL, "MEDIUM", "rotated 270deg, OCR",
                        note="OCR robustness probe — rotation."),
                finding("contact.phone", PHONE, "MEDIUM", "rotated 270deg, OCR",
                        note="OCR robustness probe — rotation."),
            ],
        },
        {
            "path": "tests/format_data/image_exif_rotated.png",
            "format": "png",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH",
                        "upright content, EXIF Orientation=6, OCR",
                        note="OCR robustness probe — EXIF-only rotation (pixel data "
                             "pre-rotated 270deg CW, no visible re-crop). Corrected by "
                             "image_extractor.py's exif_transpose step 0, before OSD/"
                             "brute-force/deskew ever run."),
                finding("contact.email", EMAIL, "MEDIUM",
                        "upright content, EXIF Orientation=6, OCR",
                        note="OCR robustness probe — EXIF-only rotation."),
                finding("contact.phone", PHONE, "MEDIUM",
                        "upright content, EXIF Orientation=6, OCR",
                        note="OCR robustness probe — EXIF-only rotation."),
            ],
        },
        {
            "path": "tests/format_data/data.csv",
            "format": "csv",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "row 1, sin column"),
                finding("contact.email", EMAIL, "MEDIUM", "row 1, email column"),
                finding("contact.phone", PHONE, "MEDIUM", "row 1, phone column"),
                finding("contact.address.postal_code", POSTAL_CODE, "MEDIUM",
                        "row 1, postal_code column"),
            ],
        },
        {
            "path": "tests/format_data/data.json",
            "format": "json",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "record 1, sin field"),
                finding("contact.email", EMAIL, "MEDIUM", "record 1, email field"),
                finding("contact.phone", PHONE, "MEDIUM", "record 1, phone field"),
                finding("contact.address.postal_code", POSTAL_CODE, "MEDIUM",
                        "record 1, note field"),
            ],
        },
        {
            "path": "tests/format_data/email_notice.eml",
            "format": "eml",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "plain-text body"),
                finding("contact.phone", PHONE, "MEDIUM", "plain-text body"),
                finding("identifier.financial.credit_card", CREDIT_CARD, "HIGH",
                        "HTML alternative body (tags stripped)"),
                finding("contact.address.postal_code", POSTAL_CODE, "MEDIUM",
                        "HTML alternative body (label 'Postal Code', tags stripped)"),
            ],
        },
        {
            "path": "tests/format_data/deck_summary.pptx",
            "format": "pptx",
            "negative_control": False,
            "expected": [
                finding("identifier.financial.sin", SIN, "HIGH", "slide 1 body text"),
                finding("contact.phone", PHONE, "MEDIUM", "slide 1 body text"),
                finding("contact.address.postal_code", POSTAL_CODE, "MEDIUM",
                        "slide 1 speaker notes (label 'Postal Code')"),
                finding("identifier.financial.credit_card", CREDIT_CARD, "HIGH",
                        "slide 2 table cell"),
            ],
        },
    ]
    return {"generated_by": "tests/make_format_data.py", "files": files}


def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    verify_planted_values()

    print("[+] Building DOCX files...")
    build_docx_full(os.path.join(OUT_DIR, "docx_full.docx"))
    build_docx_no_pii(os.path.join(OUT_DIR, "docx_no_pii.docx"))

    print("[+] Building XLSX file...")
    build_xlsx_multisheet(os.path.join(OUT_DIR, "xlsx_multisheet.xlsx"))

    print("[+] Building PDF files...")
    build_pdf_text_layer(os.path.join(OUT_DIR, "pdf_text_layer.pdf"))
    build_pdf_no_pii(os.path.join(OUT_DIR, "pdf_no_pii.pdf"))
    build_pdf_scanned_2page(os.path.join(OUT_DIR, "pdf_scanned_2page.pdf"))

    print("[+] Building PNG images...")
    build_png(os.path.join(OUT_DIR, "image_clean.png"), font_size=32)
    build_png(os.path.join(OUT_DIR, "image_small.png"), font_size=19)  # ~14pt
    build_png(os.path.join(OUT_DIR, "image_rotated.png"), font_size=32, rotate_deg=5)
    build_png(os.path.join(OUT_DIR, "image_rotated_180.png"), font_size=32, rotate_deg=180)
    build_png(os.path.join(OUT_DIR, "image_rotated_270.png"), font_size=32, rotate_deg=270)
    build_png_exif_rotated(os.path.join(OUT_DIR, "image_exif_rotated.png"), font_size=32)

    print("[+] Building CSV/JSON...")
    build_csv(os.path.join(OUT_DIR, "data.csv"))
    build_json(os.path.join(OUT_DIR, "data.json"))

    print("[+] Building EML...")
    build_eml(os.path.join(OUT_DIR, "email_notice.eml"))

    print("[+] Building PPTX...")
    build_pptx(os.path.join(OUT_DIR, "deck_summary.pptx"))

    print("[+] Writing manifest...")
    manifest = build_manifest()
    with open(MANIFEST_PATH, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2)

    total_files = len(manifest["files"])
    total_expected = sum(len(f["expected"]) for f in manifest["files"])
    print(f"[✓] Generated {total_files} files, {total_expected} expected detections")
    print(f"[✓] Manifest written to {MANIFEST_PATH}")


if __name__ == "__main__":
    main()
