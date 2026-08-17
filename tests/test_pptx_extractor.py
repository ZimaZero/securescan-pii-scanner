#!/usr/bin/env python3
"""Regression coverage for the v2.3 .pptx extractor.

Covers: title + bullet text, table cell text, speaker notes, multi-slide
order preservation, a malformed/empty slide (extraction must degrade
gracefully, never crash), a mocked per-shape failure (reported via
return_details, not silently dropped), and end-to-end detection of PII
planted in a slide body and in speaker notes via discovery.scan_path().

No real presentations are used anywhere in this file — every fixture is
synthetic, built inline with python-pptx.
"""

import json
import os
import sys
import tempfile
from pathlib import Path
from unittest.mock import patch

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation
from pptx.util import Inches

import discovery
from extractors.pptx_extractor import extract_pptx

# Reused from tests/make_format_data.py's fixed, pre-validated PII set.
SIN = "132-677-360"
CREDIT_CARD = "4111 1111 1111 1111"
PHONE = "416-555-0199"


def _title_bullets_deck() -> Presentation:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text_frame.text = "Quarterly Review"
    body = slide.placeholders[1]
    body.text_frame.text = "First bullet point"
    p2 = body.text_frame.add_paragraph()
    p2.text = "Second bullet point"
    return prs


def _table_deck() -> Presentation:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = width = height = Inches(1)
    table_shape = slide.shapes.add_table(2, 2, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "Field"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Card on file"
    table.cell(1, 1).text = CREDIT_CARD
    return prs


def _notes_deck() -> Presentation:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = width = height = Inches(1)
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.text = "Agenda"
    notes = slide.notes_slide
    notes.notes_text_frame.text = f"Reminder: caller SIN is {SIN}"
    return prs


def _multi_slide_deck() -> Presentation:
    prs = Presentation()
    for i, label in enumerate(["Alpha slide", "Beta slide", "Gamma slide"], start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = top = width = height = Inches(1)
        box = slide.shapes.add_textbox(left, top, width, height)
        box.text_frame.text = f"{label} ({i})"
    return prs


def _empty_slide_deck() -> Presentation:
    prs = Presentation()
    # Slide 1: real content.
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = width = height = Inches(1)
    box = slide1.shapes.add_textbox(left, top, width, height)
    box.text_frame.text = "Readable slide content"
    # Slide 2: genuinely empty (no shapes, no notes) — must degrade
    # (contribute nothing) without crashing extraction of the rest.
    prs.slides.add_slide(prs.slide_layouts[6])
    # Slide 3: real content again, must still be reached.
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    box3 = slide3.shapes.add_textbox(left, top, width, height)
    box3.text_frame.text = "Third slide content"
    return prs


def _end_to_end_deck() -> Presentation:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = width = height = Inches(1)
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.text = f"Contact phone on file: {PHONE}"
    notes = slide.notes_slide
    notes.notes_text_frame.text = f"Client SIN: {SIN}"
    return prs


def main() -> int:
    rows = []

    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)

        # --- Case 1: title + bullets -----------------------------------
        title_path = root / "title_bullets.pptx"
        _title_bullets_deck().save(str(title_path))
        text = extract_pptx(str(title_path))
        rows.append((
            "title + bullet text extracted in order",
            "Slide 1" in text
            and text.index("Quarterly Review") < text.index("First bullet point")
            and "Second bullet point" in text,
            repr(text),
        ))

        # --- Case 2: table cell text -------------------------------------
        table_path = root / "table.pptx"
        _table_deck().save(str(table_path))
        text = extract_pptx(str(table_path))
        rows.append((
            "table cell text extracted",
            "Field | Value" in text and f"Card on file | {CREDIT_CARD}" in text,
            repr(text),
        ))

        # --- Case 3: speaker notes -----------------------------------------
        notes_path = root / "notes.pptx"
        _notes_deck().save(str(notes_path))
        text = extract_pptx(str(notes_path))
        rows.append((
            "speaker notes extracted with label",
            "Agenda" in text and f"Notes: Reminder: caller SIN is {SIN}" in text,
            repr(text),
        ))

        # --- Case 4: multi-slide order preserved --------------------------
        multi_path = root / "multi.pptx"
        _multi_slide_deck().save(str(multi_path))
        text = extract_pptx(str(multi_path))
        rows.append((
            "multi-slide deck preserves slide order",
            text.index("Alpha slide")
            < text.index("Beta slide")
            < text.index("Gamma slide"),
            repr(text),
        ))

        # --- Case 5: empty slide degrades, never crashes -------------------
        empty_path = root / "empty_slide.pptx"
        _empty_slide_deck().save(str(empty_path))
        try:
            text = extract_pptx(str(empty_path))
            crashed = False
        except Exception as exc:  # pragma: no cover - failure path only
            text = ""
            crashed = True
            exc_detail = f"{type(exc).__name__}: {exc}"
        rows.append((
            "empty slide does not crash extraction; other slides still present",
            not crashed
            and "Readable slide content" in text
            and "Third slide content" in text,
            "no exception" if not crashed else exc_detail,
        ))

        # --- Case 6: a mocked per-shape failure is reported, not dropped --
        failure_path = root / "shape_failure.pptx"
        _title_bullets_deck().save(str(failure_path))
        with patch(
            "extractors.pptx_extractor._extract_shape_text",
            side_effect=RuntimeError("boom"),
        ):
            text, details = extract_pptx(str(failure_path), return_details=True)
        rows.append((
            "per-shape failure surfaced via pptx_parts_failed, not silently dropped",
            details.get("pptx_parts_failed", 0) >= 1
            and bool(details.get("pptx_part_failure_reasons"))
            and "Quarterly Review" not in text,
            json.dumps(details),
        ))

        # --- Case 7: metadata_only returns core properties + slide_count --
        meta = extract_pptx(str(multi_path), metadata_only=True)
        rows.append((
            "metadata_only reports slide_count",
            meta.get("slide_count") == 3,
            json.dumps(meta),
        ))

        # --- Case 8: end-to-end via discovery.scan_path() ------------------
        e2e_dir = root / "e2e"
        e2e_dir.mkdir()
        _end_to_end_deck().save(str(e2e_dir / "briefing.pptx"))
        html_report = discovery.scan_path(str(e2e_dir), verify=False, run_ner=False)
        json_report_path = os.path.splitext(html_report)[0] + ".json"
        with open(json_report_path, "r", encoding="utf-8") as fh:
            report = json.load(fh)
        file_entry = report["files"][0]
        matches = file_entry["matches"]
        found_sin = any(
            d.get("risk_level") == "HIGH"
            for d in matches.get("identifier.financial.sin", [])
        )
        found_phone = any(
            d.get("risk_level") == "MEDIUM"
            for d in matches.get("contact.phone", [])
        )
        rows.append((
            "scan_path detects SIN planted in speaker notes",
            found_sin,
            json.dumps(matches.get("identifier.financial.sin", [])),
        ))
        rows.append((
            "scan_path detects phone planted in slide body",
            found_phone,
            json.dumps(matches.get("contact.phone", [])),
        ))
        rows.append((
            "scan_path scan_status is 'scanned' (not extraction_failed)",
            file_entry.get("scan_status") == "scanned",
            file_entry.get("scan_status"),
        ))

    print(f"{'CASE':<68} {'RESULT':<7} DETAIL")
    print("-" * 110)
    for name, passed, detail in rows:
        print(f"{name:<68} {'PASS' if passed else 'FAIL':<7} {detail}")
    passed_count = sum(passed for _, passed, _ in rows)
    print("-" * 110)
    print(f"SUMMARY: {passed_count}/{len(rows)} passed")
    return 0 if passed_count == len(rows) else 1


def test_pptx_extractor():
    assert main() == 0


if __name__ == "__main__":
    raise SystemExit(main())
